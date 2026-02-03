"""Wrapper for loading pptx files"""

import asyncio
import base64
import io
from typing import IO, List, Optional

import aiofiles
import aiohttp
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture

from app.utils import load_llm

from .base import BaseDocumentLoader, ChunkType


class PowerPointLoader(BaseDocumentLoader):
    def __init__(self):
        self.presentation: Optional[PresentationType] = None
        self.filename = ""

    async def _load_pptx(self, source: str | IO[bytes]) -> None:
        logger.info(f"Loading PPTX {source}")

        if isinstance(source, str):
            if self.is_url(source):
                async with aiohttp.ClientSession() as session:
                    async with session.get(source) as response:
                        response.raise_for_status()
                        content = await response.read()
                        self.presentation = Presentation(io.BytesIO(content))
            else:
                async with aiofiles.open(source, mode="rb") as f:
                    content = await f.read()
                    self.presentation = await asyncio.to_thread(
                        Presentation, io.BytesIO(content)
                    )
        elif isinstance(source, (bytes, io.BytesIO, io.FileIO)):
            self.presentation = await asyncio.to_thread(Presentation, source)
        else:
            logger.error(f"Unsupported source type: {type(source)}")
            raise ValueError("Unsupported source type")

    def _extract_text_from_shape(self, shape: Shape) -> str:
        if shape.has_text_frame:
            if len(shape.text) > 50:
                return shape.text
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            text = " ".join(
                self._extract_text_from_shape(subshape) for subshape in shape.shapes
            )
            if len(text) > 50:
                return text
        return ""

    async def _extract_image_data(self, picture: Picture, slide_number: int) -> Document:
        """Extract and analyze a single image with LLM"""
        try:
            image = picture.image
            image_bytes = image.blob
            
            # Encode image to base64
            base64_image_data = base64.b64encode(image_bytes).decode("utf-8")
            
            # Get LLM summary (run in executor to avoid blocking)
            llm = load_llm()
            message = HumanMessage(content=[
                {
                    "type": "text",
                    "text": "What is the image about, what is in the image, and what is the image trying to convey? Explicitly mention the image in the summary and explain in detail."
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{base64_image_data}"}
                }
            ])
            
            # Run LLM call in executor to not block event loop
            response = await asyncio.get_event_loop().run_in_executor(
                None, llm.invoke, [message]
            )
            summary = response.content
            logger.info(f"Generated summary for slide {slide_number}, image: {picture.name}")

            return Document(
                page_content=summary,  # Store summary, not base64
                metadata=dict(
                    page_number=slide_number,
                    title=picture.name,
                    content_type=image.content_type,
                    width=picture.width,
                    height=picture.height,
                    type=ChunkType.image,
                    filename=self.filename,
                    image_base64=base64_image_data,  # Optional: store original
                ),
            )
        except Exception as e:
            logger.error(f"Error extracting image on slide {slide_number}: {e}")
            return None

    async def extract_data(self) -> List[Document]:
        if not self.presentation:
            raise ValueError("No presentation loaded")

        metadata = dict(
            author=self.presentation.core_properties.author,
            identifier=self.presentation.core_properties.identifier,
            created_at=self.presentation.core_properties.created,
            language=self.presentation.core_properties.language,
            modified=self.presentation.core_properties.modified,
            title=self.presentation.core_properties.title,
            subject=self.presentation.core_properties.subject,
            filename=self.filename,
        )

        # Process all slides concurrently
        async with asyncio.TaskGroup() as tg:
            tasks: List[asyncio.Task] = []
            for idx, slide in enumerate(self.presentation.slides, start=1):
                tasks.append(tg.create_task(self._process_slide(slide, idx)))

        results = []
        for task in tasks:
            slide_docs = task.result()
            # Filter out None results from failed image extractions
            results.extend([doc for doc in slide_docs if doc is not None])

        return [
            Document(
                page_content=doc.page_content, metadata={**doc.metadata, **metadata}
            )
            for doc in results
        ]

    async def _process_slide(self, slide, idx) -> List[Document]:
        """Process a single slide - extract text and images concurrently"""
        docs = []
        image_tasks = []

        # Separate text extraction (synchronous) from image extraction (async)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Create task for async image processing
                image_tasks.append(self._extract_image_data(shape, idx))
            else:
                # Extract text synchronously
                text = self._extract_text_from_shape(shape)
                if text:
                    docs.append(
                        Document(
                            page_content=text,
                            metadata={"page_number": idx, "type": ChunkType.text},
                        )
                    )

        # Process all images on this slide concurrently
        if image_tasks:
            image_results = await asyncio.gather(*image_tasks, return_exceptions=True)
            
            for result in image_results:
                if isinstance(result, Document):
                    docs.append(result)
                elif isinstance(result, Exception):
                    logger.error(f"Error processing image on slide {idx}: {result}")

        return docs

    async def load_async(self, source: str | IO[bytes]) -> List[Document]:
        self.filename = self._get_filename(source)
        await self._load_pptx(source)
        return await self.extract_data()



# """Wrapper for loading pptx files"""

# import asyncio
# import base64
# import io
# from typing import IO, List, Optional

# import aiofiles
# import aiohttp
# from langchain_core.documents import Document
# from langchain_core.messages import HumanMessage
# from loguru import logger
# from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# from pptx.presentation import Presentation as PresentationType
# from pptx.shapes.autoshape import Shape
# from pptx.shapes.picture import Picture

# from app.utils import load_llm

# from .base import BaseDocumentLoader, ChunkType


# class PowerPointLoader(BaseDocumentLoader):
#     def __init__(self):
#         self.presentation: Optional[PresentationType] = None
#         self.filename = ""

#     async def _load_pptx(self, source: str | IO[bytes]) -> None:
#         logger.info(f"Loading PPTX {source}")

#         if isinstance(source, str):
#             if self.is_url(source):
#                 async with aiohttp.ClientSession() as session:
#                     async with session.get(source) as response:
#                         response.raise_for_status()
#                         content = await response.read()
#                         self.presentation = Presentation(io.BytesIO(content))
#             else:
#                 async with aiofiles.open(source, mode="rb") as f:
#                     content = await f.read()
#                     self.presentation = await asyncio.to_thread(
#                         Presentation, io.BytesIO(content)
#                     )
#         elif isinstance(source, (bytes, io.BytesIO, io.FileIO)):
#             self.presentation = await asyncio.to_thread(Presentation, source)
#         else:
#             logger.error(f"Unsupported source type: {type(source)}")
#             raise ValueError("Unsupported source type")

#     def _extract_text_from_shape(self, shape: Shape) -> str:
#         if shape.has_text_frame:
#             if len(shape.text) > 50:
#                 return shape.text
#         elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
#             text = " ".join(
#                 self._extract_text_from_shape(subshape) for subshape in shape.shapes
#             )
#             if len(text) > 50:
#                 return text
#         return ""

#     def _extract_image_data(self, picture: Picture, slide_number: int) -> Document:
#         image = picture.image
#         image_bytes = image.blob
#         # base64_data = base64.b64encode(image_bytes).decode("utf-8")
#         base64_image_data = base64.b64encode(image_bytes).decode("utf-8")
#         llm = load_llm()
#         message = HumanMessage(content=[
#             {"type": "text", "text": "What is the image about, what is in the image, and what is the image trying to convey?, explicitly mention the image in the summary and explan in detail."},
#             {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image_data}"}}
#         ])
#         response = llm.invoke([message])
#         summary = response.content
#         logger.info(f"Summary: {summary}")

#         return Document(
#             page_content=summary,
#             metadata=dict(
#                 page_number=slide_number,
#                 title=picture.name,
#                 content_type=image.content_type,
#                 width=picture.width,
#                 height=picture.height,
#                 type=ChunkType.image,
#                 filename=self.filename,
#             ),
#         )

#     async def extract_data(self) -> List[Document]:
#         if not self.presentation:
#             raise ValueError("No presentation loaded")

#         metadata = dict(
#             author=self.presentation.core_properties.author,
#             identifier=self.presentation.core_properties.identifier,
#             created_at=self.presentation.core_properties.created,
#             language=self.presentation.core_properties.language,
#             modified=self.presentation.core_properties.modified,
#             title=self.presentation.core_properties.title,
#             subject=self.presentation.core_properties.subject,
#             filename=self.filename,
#         )

#         async with asyncio.TaskGroup() as tg:
#             tasks: List[asyncio.Task] = []
#             for idx, slide in enumerate(self.presentation.slides, start=1):
#                 tasks.append(tg.create_task(self._process_slide(slide, idx)))

#         results = []
#         for task in tasks:
#             results.extend(task.result())

#         return [
#             Document(
#                 page_content=doc.page_content, metadata={**doc.metadata, **metadata}
#             )
#             for doc in results
#         ]

#     async def _process_slide(self, slide, idx) -> List[Document]:
#         docs = []

#         for shape in slide.shapes:
#             if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#                 docs.append(self._extract_image_data(shape, idx))
#             else:
#                 text = self._extract_text_from_shape(shape)
#                 if text:
#                     docs.append(
#                         Document(
#                             page_content=text,
#                             metadata={"page_number": idx, "type": ChunkType.text},
#                         )
#                     )

#         return docs

#     async def load_async(self, source: str | IO[bytes]) -> List[Document]:
#         self.filename = self._get_filename(source)
#         await self._load_pptx(source)
#         return await self.extract_data()

